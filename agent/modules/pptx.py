"""
PowerPoint automation module — read, create, and edit .pptx files.

Contract: exposes TOOL_DEFINITIONS and dispatch() for the registry.
"""

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


TOOL_DEFINITIONS: list[dict] = [
    {
        "name": "pptx_read",
        "description": "Read all slide content from a .pptx file. Returns slide index, title, and text for each slide.",
        "input_schema": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Path to .pptx file."},
            },
            "required": ["path"],
        },
    },
    {
        "name": "pptx_create",
        "description": (
            "Create a new .pptx presentation from a list of slide definitions. "
            "Each slide has a title and optional bullet points or body text."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Output .pptx file path."},
                "title": {"type": "string", "description": "Presentation title (shown on the first title slide)."},
                "subtitle": {"type": "string", "description": "Subtitle for the title slide."},
                "slides": {
                    "type": "array",
                    "description": "List of slide definitions.",
                    "items": {
                        "type": "object",
                        "properties": {
                            "title": {"type": "string"},
                            "layout": {
                                "type": "string",
                                "enum": ["title_content", "title_only", "blank", "two_content"],
                                "description": "Slide layout. Default: title_content.",
                                "default": "title_content",
                            },
                            "bullets": {
                                "type": "array",
                                "items": {"type": "string"},
                                "description": "Bullet points for the content placeholder.",
                            },
                            "body": {"type": "string", "description": "Free-form body text (used if no bullets)."},
                        },
                        "required": ["title"],
                    },
                },
            },
            "required": ["path", "slides"],
        },
    },
    {
        "name": "pptx_add_slide",
        "description": "Add a new slide to an existing .pptx file.",
        "input_schema": {
            "type": "object",
            "properties": {
                "path": {"type": "string"},
                "title": {"type": "string"},
                "bullets": {"type": "array", "items": {"type": "string"}},
                "body": {"type": "string"},
                "layout": {
                    "type": "string",
                    "enum": ["title_content", "title_only", "blank"],
                    "default": "title_content",
                },
            },
            "required": ["path", "title"],
        },
    },
    {
        "name": "pptx_update_slide",
        "description": "Replace the title and/or body text of an existing slide by index.",
        "input_schema": {
            "type": "object",
            "properties": {
                "path": {"type": "string"},
                "slide_index": {"type": "integer", "description": "0-based slide index."},
                "title": {"type": "string", "description": "New title (optional)."},
                "body": {"type": "string", "description": "New body text (replaces all text in body placeholder)."},
            },
            "required": ["path", "slide_index"],
        },
    },
    {
        "name": "pptx_get_metadata",
        "description": "Get metadata and slide count from a .pptx file.",
        "input_schema": {
            "type": "object",
            "properties": {
                "path": {"type": "string"},
            },
            "required": ["path"],
        },
    },
]


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

_LAYOUT_MAP = {
    "title_content": 1,
    "title_only":    5,
    "blank":         6,
    "two_content":   3,
}

def _get_layout(prs: Presentation, layout: str):
    idx = _LAYOUT_MAP.get(layout, 1)
    try:
        return prs.slide_layouts[idx]
    except IndexError:
        return prs.slide_layouts[0]


def _fill_slide(slide, title: str, bullets: list[str] | None = None, body: str = "") -> None:
    if slide.shapes.title:
        slide.shapes.title.text = title

    # find the first non-title placeholder for content
    content_ph = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx != 0:
            content_ph = ph
            break

    if content_ph is None:
        return

    if bullets:
        tf = content_ph.text_frame
        tf.clear()
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = bullet
            p.level = 0
    elif body:
        content_ph.text = body


# ---------------------------------------------------------------------------
# Implementations
# ---------------------------------------------------------------------------

def pptx_read(path: str) -> dict[str, Any]:
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides):
        title = ""
        texts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape == slide.shapes.title:
                title = shape.text_frame.text
            else:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        texts.append(para.text)
        slides.append({"index": i, "title": title, "content": texts})
    return {"path": path, "slide_count": len(slides), "slides": slides}


def pptx_create(
    path: str,
    slides: list[dict],
    title: str = "",
    subtitle: str = "",
) -> dict[str, Any]:
    prs = Presentation()
    Path(path).parent.mkdir(parents=True, exist_ok=True)

    # Title slide
    if title:
        title_layout = prs.slide_layouts[0]
        ts = prs.slides.add_slide(title_layout)
        ts.shapes.title.text = title
        if subtitle and len(ts.placeholders) > 1:
            ts.placeholders[1].text = subtitle

    for slide_def in slides:
        layout_name = slide_def.get("layout", "title_content")
        sl = prs.slides.add_slide(_get_layout(prs, layout_name))
        _fill_slide(
            sl,
            slide_def.get("title", ""),
            slide_def.get("bullets"),
            slide_def.get("body", ""),
        )

    prs.save(path)
    return {"path": path, "status": "created", "slide_count": len(prs.slides)}


def pptx_add_slide(
    path: str, title: str,
    bullets: list[str] | None = None,
    body: str = "",
    layout: str = "title_content",
) -> dict[str, Any]:
    prs = Presentation(path)
    sl = prs.slides.add_slide(_get_layout(prs, layout))
    _fill_slide(sl, title, bullets, body)
    prs.save(path)
    return {"path": path, "status": "slide_added", "new_slide_index": len(prs.slides) - 1, "total_slides": len(prs.slides)}


def pptx_update_slide(
    path: str, slide_index: int,
    title: str | None = None,
    body: str | None = None,
) -> dict[str, Any]:
    prs = Presentation(path)
    if slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range (total: {len(prs.slides)}).")
    slide = prs.slides[slide_index]

    if title is not None and slide.shapes.title:
        slide.shapes.title.text = title

    if body is not None:
        for ph in slide.placeholders:
            if ph.placeholder_format.idx != 0:
                ph.text = body
                break

    prs.save(path)
    return {"path": path, "slide_index": slide_index, "status": "updated"}


def pptx_get_metadata(path: str) -> dict[str, Any]:
    prs = Presentation(path)
    core = prs.core_properties
    return {
        "path": path,
        "title": core.title or "",
        "author": core.author or "",
        "subject": core.subject or "",
        "created": str(core.created) if core.created else "",
        "modified": str(core.modified) if core.modified else "",
        "slide_count": len(prs.slides),
        "width_inches": round(prs.slide_width / 914400, 2),
        "height_inches": round(prs.slide_height / 914400, 2),
    }


def dispatch(tool_name: str, tool_input: dict) -> dict[str, Any]:
    match tool_name:
        case "pptx_read":          return pptx_read(**tool_input)
        case "pptx_create":        return pptx_create(**tool_input)
        case "pptx_add_slide":     return pptx_add_slide(**tool_input)
        case "pptx_update_slide":  return pptx_update_slide(**tool_input)
        case "pptx_get_metadata":  return pptx_get_metadata(**tool_input)
        case _: raise ValueError(f"Unknown tool: {tool_name}")
